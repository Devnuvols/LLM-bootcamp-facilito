# pip install python-docx python-pptx
# pip install beautifulsoup4 lxml requests
from docx import Document
from pptx import Presentation
from bs4 import BeautifulSoup
import requests
import re

def extract_html(html_path):
    with open(html_path, 'r', encoding='utf-8') as file:
        content = file.read()
    
    # Parse the HTML content
    soup = BeautifulSoup(content, 'lxml')
    
    # Extract all text
    text = soup.get_text(separator='\n')
    
    return text

def extract_text_from_url(url):
    # Get the HTML content from the URL
    response = requests.get(url)
    content = response.content
    
    # Parse the HTML content
    soup = BeautifulSoup(content, 'lxml')
    
    # Extract all text
    text = soup.get_text(separator='\n')
    # Replace multiple newlines with a single newline
    text = re.sub(r'\n+', '\n', text)
    return text

def extract_doc(docx_path):
    # Leer el documento Word
    doc = Document(docx_path)
    
    full_text = ""
    for paragraph in doc.paragraphs:
        full_text += paragraph.text + "\n"

    
    return full_text

def extract_text_from_pptx(pptx_path):
    # Load the presentation
    prs = Presentation(pptx_path)
    
    # Extract text from each slide
    text_runs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_runs.append(shape.text)
    
    return "\n".join(text_runs)


docx_path = 'En un lugar.docx'
text= extract_doc(docx_path)
print(text)  

pptx_path = 'Proyecto integración 2022-2.pptx'
text = extract_text_from_pptx(pptx_path)
print(text)

html_path = 'LEGO.html'
text = extract_html(html_path)
print(text)
print("--------------------")

url = 'https://www.nasa.gov/news/all-news/'
text = extract_text_from_url(url)
print(text)

